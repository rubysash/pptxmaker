import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR

# ===================== CONFIGURATION =====================
OUTPUT_PPTX = "styled_presentation.pptx"
BACKGROUND_IMAGE = "background.png"  # Set to None if no background image

# Font configurations
TITLE_FONT = {
    "size": Pt(44),
    "bold": True,
    "name": "Arial",
    "color": RGBColor(255, 0, 0)  # Red
}

SUB_HEADER_FONT = {
    "size": Pt(28),  # Slightly reduced size for better fitting
    "bold": True,
    "name": "Arial",
    "color": RGBColor(0, 0, 255)  # Blue
}

DESCRIPTION_FONT = {
    "size": Pt(22),  # Reduced size to prevent overflow
    "italic": True,
    "name": "Arial",
    "color": RGBColor(255, 255, 255)  # white
}

COLUMN_WIDTH = Inches(4)
COLUMN_HEIGHT = Inches(5)
LEFT_COLUMN_X = Inches(0.5)
RIGHT_COLUMN_X = Inches(5.5)
TEXT_Y = Inches(2)
IMAGE_Y = Inches(2)
IMAGE_WIDTH = Inches(4)
IMAGE_HEIGHT = Inches(5)
# ========================================================

def create_styled_ppt(json_file, output_pptx, background_image=None):
    # Load JSON data
    with open(json_file, 'r', encoding='utf-8') as file:
        data = json.load(file)

    # Create PowerPoint Presentation
    prs = Presentation()

    for header, slide_data in data.items():
        alignment = slide_data.get("alignment", "left")
        image_path = slide_data.get("image", None)
        content = slide_data.get("content", {})

        # Add a slide
        slide_layout = prs.slide_layouts[5]  # Blank slide for full customization
        slide = prs.slides.add_slide(slide_layout)

        # Set background image if provided
        if background_image:
            slide.shapes.add_picture(background_image, 0, 0, prs.slide_width, prs.slide_height)

        # Add Title (Centered)
        title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        title_frame = title_shape.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = header
        title_p.font.size = TITLE_FONT["size"]
        title_p.font.bold = TITLE_FONT["bold"]
        title_p.font.name = TITLE_FONT["name"]
        title_p.font.color.rgb = TITLE_FONT["color"]
        title_p.alignment = PP_ALIGN.CENTER

        # Determine column positions based on alignment
        text_x = LEFT_COLUMN_X if alignment == "left" else RIGHT_COLUMN_X
        image_x = RIGHT_COLUMN_X if alignment == "left" else LEFT_COLUMN_X

        # Add Content Box
        content_shape = slide.shapes.add_textbox(text_x, TEXT_Y, COLUMN_WIDTH, COLUMN_HEIGHT)
        content_frame = content_shape.text_frame
        content_frame.word_wrap = True
        content_frame.auto_size = None
        content_frame.margin_top = Pt(2)
        content_frame.margin_bottom = Pt(2)
        content_frame.margin_left = Pt(4)
        content_frame.margin_right = Pt(4)
        content_frame.vertical_anchor = MSO_ANCHOR.TOP  # Ensures text starts from the top

        for sub_header, description in content.items():
            # Add Sub-header
            p = content_frame.add_paragraph()
            p.text = sub_header
            p.font.size = SUB_HEADER_FONT["size"]
            p.font.bold = SUB_HEADER_FONT["bold"]
            p.font.name = SUB_HEADER_FONT["name"]
            p.font.color.rgb = SUB_HEADER_FONT["color"]
            p.alignment = PP_ALIGN.LEFT

            # Add Description
            p_desc = content_frame.add_paragraph()
            p_desc.text = description
            p_desc.font.size = DESCRIPTION_FONT["size"]
            p_desc.font.italic = DESCRIPTION_FONT["italic"]
            p_desc.font.name = DESCRIPTION_FONT["name"]
            p_desc.font.color.rgb = DESCRIPTION_FONT["color"]
            p_desc.space_after = Pt(12)  # Reduce spacing to prevent overflow
            p_desc.alignment = PP_ALIGN.LEFT

        # Add Image on the opposite column
        if image_path:
            try:
                slide.shapes.add_picture(image_path, image_x, IMAGE_Y, IMAGE_WIDTH, IMAGE_HEIGHT)
            except Exception as e:
                print(f"Warning: Could not add image '{image_path}' - {e}")

    # Save the PowerPoint file
    prs.save(output_pptx)
    print(f"PPTX file '{output_pptx}' created successfully.")

# Example Usage
json_file = "slides_data.json"  # JSON input file
create_styled_ppt(json_file, OUTPUT_PPTX, BACKGROUND_IMAGE)
